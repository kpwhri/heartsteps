import os
import uuid
import datetime
import configparser
import pymongo
import psycopg2
import urllib
import pandas as pd
import numpy as np
import subprocess
from pymongo import UpdateOne
from pptx import Presentation
import matplotlib.pyplot as plt
import seaborn as sns
from matplotlib.patches import Patch

def get_config_obj():
    config = configparser.ConfigParser()
    config.read('justwalk_secrets/db_dump_config.ini')
    
    return config

def get_mongodb_client():
    config = get_config_obj()
    mongodb_config = config['MongoDB']
    mongodb_uri = "mongodb://{}:{}@{}:{}/{}".format(mongodb_config['username'], urllib.parse.quote(mongodb_config['password']), mongodb_config['host'], mongodb_config['port'], mongodb_config['db'])
    client = pymongo.MongoClient(mongodb_uri)
    return client

def get_mongodb_db(dbname='justwalk'):
    client = get_mongodb_client()
    db = client[dbname]
    return db

def get_mongodb_collection(collection_name):
    db = get_mongodb_db()
    collection = db[collection_name]
    return collection

def get_postgres_connection():
    """
    Get a connection to PostgreSQL
    """
    config = get_config_obj()
    postgres_config = config['PostgreSQL']
    conn = psycopg2.connect(
        host=postgres_config['host'],
        port=postgres_config['port'],
        database=postgres_config['db'],
        user=postgres_config['username'],
        password=postgres_config['password']
    )
    return conn

def dump_table(table_name, columns, where_clause=None) -> pd.DataFrame:
    """
    Dump a table from PostgreSQL to a pandas DataFrame
    """
    with get_postgres_connection() as conn:
        with conn.cursor() as cur:
            if where_clause is not None:
                cur.execute("SELECT {} FROM {} WHERE {}".format(','.join(columns), table_name, where_clause))
            else:
                cur.execute("SELECT {} FROM {}".format(','.join(columns), table_name))
            rows = cur.fetchall()
            df = pd.DataFrame(rows, columns=columns)
    return df

def dump_to_collection(df, collection_name, key_field='id'):
    """
    Dump a pandas DataFrame to a MongoDB collection
    """
    collection = get_mongodb_collection(collection_name)

    if df.shape[0] > 0:
        for index, column in enumerate(df.columns):
            if df[column].dtype == 'datetime64[ns]':
                df[column] = df[column].astype('str')
            elif isinstance(df[column][0], datetime.date):
                df[column] = df[column].astype('str')
            elif isinstance(df[column][0], datetime.datetime):
                df[column] = df[column].astype('str')

        if isinstance(key_field, str):
            return collection.bulk_write([
                UpdateOne(
                    {key_field: row[key_field]},
                    {'$set': row.to_dict()},
                    upsert=True
                ) for _, row in df.iterrows()])
        elif isinstance(key_field, list):
            bulk_write_buffer = []
            for _, row in df.iterrows():
                filter_dict = {}
                for key in key_field:
                    filter_dict[key] = row[key]

                bulk_write_buffer.append(
                    UpdateOne(
                        filter_dict,
                        {'$set': row.to_dict()},
                        upsert=True
                    )
                )
            if len(bulk_write_buffer) > 0:
                return collection.bulk_write(bulk_write_buffer)
            else:
                return None
        else:
            raise ValueError('key_field must be a string or a list of strings')
    else:
        return None

def drop_all_collections(db):
    """
    Drop all collections in the database
    :param db: the database instance
    """
    for collection_name in db.list_collection_names():
        db.drop_collection(collection_name)

def filter_collection(db, tdb, collection_name, filter_dict, projection_dict, refresh_force=False) -> list:
    """
    Filter the collection and insert the filtered data into the transformed database.
    :param db: the database instance
    :param tdb: the transformed database instance
    :param collection_name: the collection name
    :param filter_dict: the filter dictionary
    :param projection_dict: the projection dictionary
    :return: the filtered data as a list of dictionaries
    """
    collection_source = db[collection_name]
    df = pd.DataFrame(
        list(collection_source.find(filter_dict, projection_dict)))
    collection_target = tdb[collection_name]

    # insert the data into the collection only if the collection's document count is not the same as the dataframe's row count
    if refresh_force or collection_target.count_documents({}) != df.shape[0]:
        collection_target.delete_many({})
        df_dict = df.to_dict('records')
        collection_target.insert_many(df_dict, ordered=False)
        return df_dict
    else:
        return list(collection_target.find({}, projection_dict))


def build_df_from_collection(db, collection_name, filter_dict, projection_dict, rename_columns=None) -> pd.DataFrame:
    """
    Build a pandas DataFrame from a MongoDB collection
    :param db: the database instance
    :param collection_name: the collection name
    :param filter_dict: the filter dictionary
    :param projection_dict: the projection dictionary
    :return: the pandas DataFrame
    """
    collection = db[collection_name]
    df = pd.DataFrame(
        list(collection.find(filter_dict, projection_dict)))
    
    # rename the columns
    if rename_columns is not None:
        df.rename(columns=rename_columns, inplace=True)
    return df

def extend_df_with_collection(df, db, collection_name, filter_dict, projection_dict, on, how='right', rename_columns:list=None) -> pd.DataFrame:
    """
    Extend a pandas DataFrame with a MongoDB collection
    :param df: the pandas DataFrame
    :param db: the database instance
    :param collection_name: the collection name
    :param filter_dict: the filter dictionary
    :param projection_dict: the projection dictionary
    :param on: the join key
    :param how: the join type
    :param rename_columns: the dictionary of columns to rename
    :return: the extended pandas DataFrame
    """
    collection = db[collection_name]
    df2 = pd.DataFrame(
        list(collection.find(filter_dict, projection_dict)))
    
    # rename the columns
    if rename_columns is not None:
        df2.rename(columns=rename_columns, inplace=True)
    
    df = pd.merge(df, df2, on=on, how=how)
    return df




def get_uuid4_filename(filename, output_dir=None):
    """
    Get a UUID4 filename
    :param filename: the original filename
    :return: the UUID4 filename
    """
    filename, ext = os.path.splitext(filename)

    if output_dir is not None:
        return os.path.join(output_dir, '{}_{}{}'.format(filename, uuid.uuid4(), ext))
    else:
        return '{}_{}{}'.format(filename, uuid.uuid4(), ext)

def get_intervention_daily_df(original):
    """
    Get the intervention daily dataframe
    
    :param original: the original dataframe
    :return: the intervention daily dataframe
    """
    new_df = original.query('day_index >= 10 and day_index <= 252').copy()
    new_df['day_index'] = new_df['day_index'].astype(int, errors='ignore')
    new_df['level_int'] = new_df['level_int'].astype(int, errors='ignore')
    return new_df

def draw_heatmap(df, index, columns, values, legend_labels, xlabel, ylabel, legend_title, figure_name, output_dir, cm_name='coolwarm') -> str:
    """
    Draw a heatmap of the given dataframe

    :param df: dataframe
    :param index: index column name
    :param columns: columns column name
    :param values: values column name
    :param legend_labels: legend labels
    :param xlabel: x label
    :param ylabel: y label
    :param legend_title: legend title
    :param figure_name: figure name
    :param output_dir: output directory
    :param cm_name: colormap name
    :return: figure path
    """
    # Create a figure
    plt.figure(figsize=(10, 5))

    # Colormap settings
    cm = plt.cm.get_cmap(cm_name)

    # Create a pivot table
    heatmap_data = df.pivot_table(index=index, columns=columns, values=values, aggfunc=np.sum)

    # Draw a heatmap with the numeric values in each cell
    if legend_labels is None:
        ax = sns.heatmap(heatmap_data, cmap=cm_name, cbar=True)
    else:
        ax = sns.heatmap(heatmap_data, cmap=cm_name, cbar=False)

    # Add title, legend and labels
    plt.xlabel(xlabel)
    plt.ylabel(ylabel)

    if legend_labels is not None:
        # Manually set the level labels
        legend_colors = [cm(x/(len(legend_labels)-1)) for x in range(len(legend_labels))]
        legend_patches = [Patch(color=color, label=label) for color, label in zip(legend_colors, legend_labels)]
        ax.legend(handles=legend_patches, loc='lower center', bbox_to_anchor=(0.5, -0.3), ncol=len(legend_labels), title=legend_title)
    
    # Save heatmap to file
    figure_filepath = get_uuid4_filename(figure_name, output_dir)
    plt.savefig(figure_filepath, bbox_inches='tight', dpi=200)

    # Close the figure
    plt.close()
    return figure_filepath

class JWPresentation:
    """
    JustWalk Presentation
    """
    def __init__(self):
        """
        Initialize the presentation
        """
        self.prs = Presentation()
        self.title_slide_layout = self.prs.slide_layouts[0]
        self.contents_slide_layout = self.prs.slide_layouts[1]
        self.__add_title_slide()
        self.toc = {}
        
    def __add_title_slide(self):
        """
        Add the title slide
        """
        title_text = "JustWalk JITAI Data Visualization"
        subtitle_text = "Junghwan Park\n{}".format(pd.Timestamp.today().strftime('%Y-%m-%d'))

        slide = self.prs.slides.add_slide(self.title_slide_layout)
        title = slide.placeholders[0]
        subtitle = slide.placeholders[1]

        title.text = title_text
        subtitle.text = subtitle_text

    def build(self):
        """
        Build the presentation
        """
        self.__section_slide(self.toc)

    def __section_slide(self, toc_item):
        """
        Add a section slide

        :param toc_item: the toc item
        """
        slide = self.prs.slides.add_slide(self.contents_slide_layout)
        title = slide.shapes.title
        title.text = toc_item['title']

        item_texts = [item['title'] for item in toc_item['items']]
        contents = slide.placeholders[1]
        tf = contents.text_frame
        tf.text = '\n'.join(item_texts)

        self.__add_slides(toc_item['items'])

    def __add_slides(self, toc_items):
        """
        Add slides

        :param toc_items: the toc items
        """
        for toc_item in toc_items:
            if toc_item['type'] == 'section':
                self.__section_slide(toc_item)
            elif toc_item['type'] == 'slide':
                self.__normal_slide(toc_item)
            else:
                raise ValueError(
                    'Invalid toc_item type: {}'.format(toc_item['type']))

    def __normal_slide(self, toc_item):
        """
        Add a normal slide

        :param toc_item: the toc item
        """
        contents_slide_layout = self.prs.slide_layouts[1]
        title_only_slide_layout = self.prs.slide_layouts[5]
        
        if 'figure' in toc_item:
            slide = self.prs.slides.add_slide(contents_slide_layout)
            title = slide.shapes.title
            title.text = toc_item['title']

            image_path = toc_item['figure']

            # remove the text box
            for shape in slide.shapes:
                if shape.has_text_frame and shape != title:
                    shape.text_frame.clear()

            # insert the image to the slide
            left = 0.05
            top = 0.2
            right = left
            
            width_ = 1 - (left + right)
            
            left_ = self.prs.slide_width * left
            top_ = self.prs.slide_height * top
            width_ = self.prs.slide_width * width_
            
            pic = slide.shapes.add_picture(image_path, left_, top_, width=width_)
        else:
            slide = self.prs.slides.add_slide(contents_slide_layout)
            title = slide.shapes.title
            title.text = toc_item['title']

    def save(self, filepath):
        """
        Save the presentation

        :param filepath: the file path
        """
        self.filepath = filepath
        self.prs.save(filepath)
    
    def open(self):
        """
        Open the presentation
        """
        if not hasattr(self, "filepath"):
            self.save('/tmp/temp.pptx')
        subprocess.call(['open', self.filepath])