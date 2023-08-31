from config import SETTINGS_OUTPUT_DIR, SETTINGS_OUTPUT_FILENAME, MONGO_DB_URI_DESTINATION
from utils import get_database, get_intervention_daily_df, draw_heatmap, draw_distribution_heatmap, JWPresentation, JWSection, draw_sorted_bars
import os
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np
from matplotlib.patches import Patch
from pptx import Presentation
import logging

def prepare():
    # create output directory if it does not exist
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)

def form_the_presentation(filename_dict):
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)
    filepath = os.path.join(output_dir, "{} ({}).pptx".format(SETTINGS_OUTPUT_FILENAME, pd.Timestamp.today().strftime('%Y-%m-%d')))

    jwp = JWPresentation()

    # add a table of contents
    toc = JWSection('Table of Contents')
    daily_section = toc.add_section('Daily Data')
    daily_section.add_slide('Levels', filename_dict['levels'])

    goals_section = daily_section.add_section('Goals')
    goals_section.add_slide('Goals', filename_dict['goals'])
    goals_section.add_slide('Goals Distribution', filename_dict['goals_distribution'])

    steps_section = daily_section.add_section('Steps')
    steps_section.add_slide('Steps', filename_dict['steps'], note='* Note: steps are capped at 20,000')
    steps_section.add_slide('Steps Distribution', filename_dict['steps_distribution'], note='* Note: steps are capped at 20,000')

    heart_rates_section = daily_section.add_section('Heart Rates')
    heart_rates_section.add_slide('Wear Time Percentage', filename_dict['wearing_time'])
    heart_rates_section.add_slide('# of Days with >60% Wear Time (Sorted)', filename_dict['wearing_time_sorted_bars'])
    
    survey_section = daily_section.add_section('Survey')
    survey_section.add_slide('Daily EMA', filename_dict['survey_daily_ema'])
    survey_section.add_slide('# of Days with answered daily EMA (Sorted)', filename_dict['survey_daily_ema_sorted_bars'])

    jwp.toc = toc.to_dict()
    logging.debug(jwp.toc)
    
    jwp.build()
    jwp.save(filepath)
    jwp.open()



def draw_level_heatmap():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    daily = pd.DataFrame(list(db['daily'].find()))

    # 1. draw a heatmap of the levels
    level_df = get_intervention_daily_df(daily)
    figure_levels = draw_heatmap(level_df, 
                                index='day_index', 
                                columns='user_id', 
                                values='level_int', 
                                legend_labels=[
                                    'Random (RA)', 
                                    'N+R (NR)', 
                                    'N+O (NO)', 
                                    'Full (FU)'
                                    ], 
                                xlabel='User ID', 
                                ylabel='Day Index', 
                                legend_title='Values', 
                                figure_name='levels.png', 
                                output_dir=output_dir)

    return figure_levels

def draw_goal_heatmap():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    daily = pd.DataFrame(list(db['daily'].find()))

    # 2. draw a heatmap of the goals
    goals_df = get_intervention_daily_df(daily)
    figure_goals = draw_heatmap(goals_df, 
                                index='day_index', 
                                columns='user_id', 
                                values='step_goal', 
                                legend_labels=None, 
                                xlabel='User ID', 
                                ylabel='Day Index', 
                                legend_title='Step Goals', 
                                figure_name='goals.png', 
                                output_dir=output_dir)

    return figure_goals

def draw_goal_distribution_heatmap():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    daily = pd.DataFrame(list(db['daily'].find()))

    # 2. draw a heatmap of the goals
    goals_df = get_intervention_daily_df(daily)

    # 3. draw a heatmap of the goals distribution
    figure_path = draw_distribution_heatmap(
                                goals_df, 
                                index='user_id', 
                                values='step_goal', 
                                xlabel='User ID', 
                                ylabel='Goals', 
                                legend_title='Step Goals', 
                                figure_name='goals_distribution.png', 
                                n=20,
                                output_dir=output_dir)

    return figure_path

def draw_steps_heatmap():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    daily = pd.DataFrame(list(db['daily'].find()))

    # 3. draw a heatmap of the steps
    steps_df = get_intervention_daily_df(daily)

    # 4. cut off the steps at 20,000
    steps_df.loc[steps_df['steps'] > 20000, 'steps'] = 20000

    figure_steps = draw_heatmap(steps_df, 
                                index='day_index', 
                                columns='user_id', 
                                values='steps', 
                                legend_labels=None, 
                                xlabel='User ID', 
                                ylabel='Day Index', 
                                legend_title='Steps', 
                                figure_name='steps.png', 
                                output_dir=output_dir)

    return figure_steps

def draw_steps_distribution_heatmap():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    daily = pd.DataFrame(list(db['daily'].find()))

    # 3. draw a heatmap of the steps
    steps_df = get_intervention_daily_df(daily)

    # 4. cut off the steps at 20,000
    steps_df.loc[steps_df['steps'] > 20000, 'steps'] = 20000

    # 5. draw a heatmap of the steps distribution
    figure_path = draw_distribution_heatmap(
                                steps_df, 
                                index='user_id', 
                                values='steps', 
                                xlabel='User ID', 
                                ylabel='Steps', 
                                legend_title='Steps', 
                                figure_name='steps_distribution.png', 
                                n=20,
                                output_dir=output_dir)

    return figure_path


def draw_wearing_time_heatmap():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    daily = pd.DataFrame(list(db['daily'].find()))

    # draw a heatmap of the wearing time percentage
    wearing_time_pct_df = get_intervention_daily_df(daily)


    figure_wearing_time_pct = draw_heatmap(wearing_time_pct_df, 
                                index='day_index', 
                                columns='user_id', 
                                values='wearing_minutes', 
                                legend_labels=None, 
                                xlabel='User ID', 
                                ylabel='Day Index', 
                                legend_title='Wearing Minutes', 
                                figure_name='wearing_time.png', 
                                output_dir=output_dir)

    return figure_wearing_time_pct

def draw_wearing_time_sorted_bars():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    daily = pd.DataFrame(list(db['daily'].find()))

    # draw a heatmap of the wearing time percentage
    wearing_time_pct_df = get_intervention_daily_df(daily)

    # create a data frame with the count of days with over 60% wearing time
    agg_wearing_time_pct_df = wearing_time_pct_df.groupby('user_id').agg({'wearing_minutes': lambda x: (x > 0.6*1440).sum()})
    agg_wearing_time_pct_df = agg_wearing_time_pct_df.reset_index()

    # sort the data frame by the count of days with over 60% wearing time
    agg_wearing_time_pct_df = agg_wearing_time_pct_df.sort_values(by='wearing_minutes', ascending=False)

    # rename the column
    agg_wearing_time_pct_df = agg_wearing_time_pct_df.rename(columns={'wearing_minutes': 'days_over60'})
    
    # 6. draw a bar chart of the wearing time percentage
    figure_wearing_time_sorted = draw_sorted_bars(agg_wearing_time_pct_df, 
                                index='user_id', 
                                values='days_over60', 
                                xlabel='User ID (Ranked)', 
                                ylabel='Number of Days with >60% Wearing Time', 
                                legend_title='Wearing Time Percentage', 
                                figure_name='wearing_time_sorted.png', 
                                output_dir=output_dir)

    return figure_wearing_time_sorted

def draw_survey_daily_ema_heatmap():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    df = pd.DataFrame(list(db['daily'].find()))

    # draw a heatmap of the daily ema
    daily_ema_df = get_intervention_daily_df(df)
    daily_ema_df['daily_ema_answered'] = [1 if x == True else 0 for x in daily_ema_df['daily_ema_answered']]

    # 7. draw a heatmap of the daily ema
    figure_daily_ema = draw_heatmap(daily_ema_df, 
                                index='day_index', 
                                columns='user_id', 
                                values='daily_ema_answered', 
                                legend_labels=['Not Answered', 'Answered'], 
                                xlabel='User ID', 
                                ylabel='Day Index', 
                                legend_title='Daily EMA', 
                                figure_name='daily_ema.png', 
                                output_dir=output_dir)

    return figure_daily_ema

def draw_survey_daily_ema_sorted_bars():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    df = pd.DataFrame(list(db['daily'].find()))

    # draw a heatmap of the daily ema
    daily_ema_df = get_intervention_daily_df(df)
    
    # create a data frame with the count of days with over 60% wearing time
    agg_daily_ema_df = daily_ema_df.groupby('user_id').agg({'daily_ema_answered': lambda x: (x == True).sum()})
    agg_daily_ema_df = agg_daily_ema_df.reset_index()

    # sort the data frame by the count of days with over 60% wearing time
    agg_daily_ema_df = agg_daily_ema_df.sort_values(by='daily_ema_answered', ascending=False)

    # rename the column
    agg_daily_ema_df = agg_daily_ema_df.rename(columns={'daily_ema_answered': 'days_answered'})
    
    # 7. draw a bar chart of the daily ema
    figure_daily_ema_sorted = draw_sorted_bars(agg_daily_ema_df, 
                                index='user_id', 
                                values='days_answered', 
                                xlabel='User ID (Ranked)', 
                                ylabel='Number of Days with Daily EMA Answered', 
                                legend_title='Daily EMA', 
                                figure_name='daily_ema_sorted.png', 
                                output_dir=output_dir)

    return figure_daily_ema_sorted

def draw_fitbit_update_heatmap():
    # Default Visualization Parameters
    cm_name = 'coolwarm'
    output_dir = os.path.join(os.getcwd(), SETTINGS_OUTPUT_DIR)

    # get the database
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    df = pd.DataFrame(list(db['daily'].find()))

    # draw a heatmap of the daily ema
    daily_ema_df = get_intervention_daily_df(df)
    
    # create a data frame with the count of days with over 60% wearing time
    agg_daily_ema_df = daily_ema_df.groupby('user_id').agg({'daily_ema_answered': lambda x: (x == True).sum()})
    agg_daily_ema_df = agg_daily_ema_df.reset_index()

    # sort the data frame by the count of days with over 60% wearing time
    agg_daily_ema_df = agg_daily_ema_df.sort_values(by='daily_ema_answered', ascending=False)

    # rename the column
    agg_daily_ema_df = agg_daily_ema_df.rename(columns={'daily_ema_answered': 'days_answered'})
    
    # 7. draw a bar chart of the daily ema
    figure_daily_ema_sorted = draw_sorted_bars(agg_daily_ema_df, 
                                index='user_id', 
                                values='days_answered', 
                                xlabel='User ID (Ranked)', 
                                ylabel='Number of Days with Daily EMA Answered', 
                                legend_title='Daily EMA', 
                                figure_name='daily_ema_sorted.png', 
                                output_dir=output_dir)

    return figure_daily_ema_sorted