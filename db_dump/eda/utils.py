from config import MONGO_DB_URI_DESTINATION
from pymongo import MongoClient
import pandas as pd
import numpy as np
import seaborn as sns
import matplotlib.pyplot as plt
import os
import uuid
import datetime
import subprocess
from matplotlib.patches import Patch
from pptx import Presentation

def get_database(uri:str, database_name:str):
    client = MongoClient(uri)
    return client[database_name]

def get_participant_list():
    db = get_database(MONGO_DB_URI_DESTINATION, 'justwalk')
    participants = db['participants']
    return list(participants.distinct(key='user_id'))

def get_intervention_daily_df(original):
    """
    Get the intervention daily dataframe
    
    :param original: the original dataframe
    :return: the intervention daily dataframe
    """
    new_df = original.query('day_index >= 10 and day_index <= 252').copy()
    new_df['day_index'] = new_df['day_index'].astype(int, errors='ignore')
    new_df['level_int'] = new_df['level_int'].astype(int, errors='ignore')
    return new_df


def draw_heatmap(df, index, columns, values, legend_labels, xlabel, ylabel, legend_title, figure_name, output_dir, cm_name='coolwarm') -> str:
    """
    Draw a heatmap of the given dataframe

    :param df: dataframe
    :param index: index column name
    :param columns: columns column name
    :param values: values column name
    :param legend_labels: legend labels
    :param xlabel: x label
    :param ylabel: y label
    :param legend_title: legend title
    :param figure_name: figure name
    :param output_dir: output directory
    :param cm_name: colormap name
    :return: figure path
    """
    # Create a figure
    plt.figure(figsize=(10, 5))

    # Colormap settings
    cm = plt.cm.get_cmap(cm_name)

    # Create a pivot table
    heatmap_data = df.pivot_table(index=index, columns=columns, values=values, aggfunc=np.sum)

    # Draw a heatmap with the numeric values in each cell
    if legend_labels is None:
        ax = sns.heatmap(heatmap_data, cmap=cm_name, cbar=True)
    else:
        ax = sns.heatmap(heatmap_data, cmap=cm_name, cbar=False)

    # Add title, legend and labels
    plt.xlabel(xlabel)
    plt.ylabel(ylabel)

    if legend_labels is not None:
        # Manually set the level labels
        legend_colors = [cm(x/(len(legend_labels)-1)) for x in range(len(legend_labels))]
        legend_patches = [Patch(color=color, label=label) for color, label in zip(legend_colors, legend_labels)]
        ax.legend(handles=legend_patches, loc='lower center', bbox_to_anchor=(0.5, -0.3), ncol=len(legend_labels), title=legend_title)
    
    # Save heatmap to file
    figure_filepath = get_uuid4_filename(figure_name, output_dir)
    plt.savefig(figure_filepath, bbox_inches='tight', dpi=200)

    # Close the figure
    plt.close()
    return figure_filepath

def draw_distribution_heatmap(df, index, values, xlabel, ylabel, legend_title, figure_name, n, output_dir, cm_name='coolwarm') -> str:
    """
    Draw a heatmap of the given dataframe

    :param df: dataframe
    :param index: index column name
    :param values: values column name
    :param xlabel: x label
    :param ylabel: y label
    :param legend_title: legend title
    :param figure_name: figure name
    :param n: number of levels
    :param output_dir: output directory
    :param cm_name: colormap name
    :return: figure path
    """
    # Create a figure
    plt.figure(figsize=(10, 5))

    # Colormap settings
    cm = plt.cm.get_cmap(cm_name)

    # Get the maximum and minimum values
    max_value = df[values].max()
    min_value = df[values].min()

    # Get the bins of the values with the count of N
    bins = np.linspace(min_value, max_value, n+1)

    # Get the count in each bin for each user and convert it to a pd.DataFrame with three columns (index, bin_median, count) to be used in the heatmap
    bins_df = df.groupby([index, pd.cut(df[values], bins)]).size().reset_index().rename(columns={0: 'count'})

    # Convert the bins_df to a pd.DataFrame with pivot_table
    bins_df = bins_df.pivot_table(index=values, columns=index, values='count', aggfunc=np.sum)

    # Draw the heatmap for bins count for each user. Rows are bins and columns are users
    ax = sns.heatmap(bins_df, cmap=cm_name, cbar=False)

    # Add title, legend and labels
    plt.xlabel(xlabel)
    plt.ylabel(ylabel)

    # Save heatmap to file
    figure_filepath = get_uuid4_filename(figure_name, output_dir)
    plt.savefig(figure_filepath, bbox_inches='tight', dpi=200)

    # Close the figure
    plt.close()
    return figure_filepath

def draw_scatterplot(df, x, y, xlabel, ylabel, figure_name, output_dir, regline=False, x_jitter=0, y_jitter=0, size=0.5) -> str:
    """
    Draw a scatterplot of the given dataframe

    :param df: dataframe
    :param x: x column name
    :param y: y column name
    :param xlabel: x label
    :param ylabel: y label
    :param legend_title: legend title
    :param figure_name: figure name
    :param output_dir: output directory
    :return: figure path
    """
    # Create a figure
    plt.figure(figsize=(10, 5))

    if regline:
        # Draw a scatterplot with a jitter and a regression line
        ax = sns.regplot(x=x, y=y, data=df, x_jitter=x_jitter, y_jitter=y_jitter)
    else:
        def jitter(x, jitter_amount):
            return x + np.random.uniform(-jitter_amount, jitter_amount, len(x))
        nx = jitter(df[x], x_jitter)
        ny = jitter(df[y], y_jitter)
        # draw a scatterplot with a jitter, and a size of 0.1
        ax = sns.scatterplot(x=nx, y=ny, data=df, s=size)
    
    # Add title, legend and labels
    plt.xlabel(xlabel)
    plt.ylabel(ylabel)

    # Save heatmap to file
    figure_filepath = get_uuid4_filename(figure_name, output_dir)
    plt.savefig(figure_filepath, bbox_inches='tight', dpi=200)

    # Close the figure
    plt.close()
    return figure_filepath


def draw_sorted_bars(df, 
                    index, 
                    values, 
                    xlabel, 
                    ylabel, 
                    legend_title, 
                    figure_name, 
                    output_dir) -> str:
    """
    Draw a sorted bar plot of the given dataframe

    :param df: dataframe
    :param index: index column name
    :param values: values column name
    :param xlabel: x label
    :param ylabel: y label
    :param legend_title: legend title
    :param figure_name: figure name
    :param output_dir: output directory
    :return: figure path
    """

    # Create a figure
    plt.figure(figsize=(10, 5))

    # Sort the dataframe by the values
    df = df.sort_values(by=values, ascending=False)
    df['ranking'] = range(1, len(df)+1)
    
    # Draw a bar plot of pd.Series
    ax = sns.barplot(x='ranking', y=values, data=df)
    ax.set_xticklabels(df[index])

    # Add title, legend and labels
    plt.xlabel(xlabel)
    plt.ylabel(ylabel)

    # rotate x labels 90 degrees
    plt.xticks(rotation=90)
    
    # Save heatmap to file
    figure_filepath = get_uuid4_filename(figure_name, output_dir)
    plt.savefig(figure_filepath, bbox_inches='tight', dpi=200)

    # Close the figure
    plt.close()
    return figure_filepath



def get_uuid4_filename(filename, output_dir=None):
    """
    Get a UUID4 filename
    :param filename: the original filename
    :return: the UUID4 filename
    """
    filename, ext = os.path.splitext(filename)

    if output_dir is not None:
        return os.path.join(output_dir, '{}_{}{}'.format(filename, uuid.uuid4(), ext))
    else:
        return '{}_{}{}'.format(filename, uuid.uuid4(), ext)

class JWSection:
    """
    JustWalk Section
    """
    def __init__(self, title):
        """
        Initialize the section
        """
        self.title = title
        self.items = []
    
    def add_section(self, title) -> 'JWSection':
        """
        Add a section

        :param title: the section title
        """
        section = JWSection(title)
        self.items.append(section)
        return section

    def add_slide(self, title, figure, note=None) -> dict:
        """
        Add a slide

        :param title: the slide title
        :param figure: the slide figure
        """
        slide_dict = {
            'title': title,
            'type': 'slide',
            'figure': figure
        }
        if note is not None:
            slide_dict['note'] = note

        self.items.append(slide_dict)

        return slide_dict
    
    def to_dict(self) -> dict:
        """
        Convert the section to a dictionary
        """
        item_dicts = []
        for item in self.items:
            if isinstance(item, JWSection):
                item_dicts.append(item.to_dict())
            else:
                item_dicts.append(item)
        
        return {
            'title': self.title,
            'type': 'section',
            'items': item_dicts
        }

class JWPresentation:
    """
    JustWalk Presentation
    """
    def __init__(self):
        """
        Initialize the presentation
        """
        self.prs = Presentation()
        self.title_slide_layout = self.prs.slide_layouts[0]
        self.contents_slide_layout = self.prs.slide_layouts[1]
        self.__add_title_slide()
        self.toc = {}
        
    def __add_title_slide(self):
        """
        Add the title slide
        """
        title_text = "JustWalk JITAI Data Visualization"
        subtitle_text = "Junghwan Park\n{}".format(pd.Timestamp.today().strftime('%Y-%m-%d'))

        slide = self.prs.slides.add_slide(self.title_slide_layout)
        title = slide.placeholders[0]
        subtitle = slide.placeholders[1]

        title.text = title_text
        subtitle.text = subtitle_text

    def build(self):
        """
        Build the presentation
        """
        self.__section_slide(self.toc)

    def __section_slide(self, toc_item):
        """
        Add a section slide

        :param toc_item: the toc item
        """
        slide = self.prs.slides.add_slide(self.contents_slide_layout)
        title = slide.shapes.title
        title.text = toc_item['title']

        item_texts = [item['title'] for item in toc_item['items']]
        contents = slide.placeholders[1]
        tf = contents.text_frame
        tf.text = '\n'.join(item_texts)

        self.__add_slides(toc_item['items'])

    def __add_slides(self, toc_items):
        """
        Add slides

        :param toc_items: the toc items
        """
        for toc_item in toc_items:
            if toc_item['type'] == 'section':
                self.__section_slide(toc_item)
            elif toc_item['type'] == 'slide':
                self.__normal_slide(toc_item)
            else:
                raise ValueError(
                    'Invalid toc_item type: {}'.format(toc_item['type']))

    def __normal_slide(self, toc_item):
        """
        Add a normal slide

        :param toc_item: the toc item
        """
        contents_slide_layout = self.prs.slide_layouts[1]
        title_only_slide_layout = self.prs.slide_layouts[5]
        
        if 'figure' in toc_item:
            slide = self.prs.slides.add_slide(contents_slide_layout)
            title = slide.shapes.title
            title.text = toc_item['title']

            image_path = toc_item['figure']

            # remove the text box
            for shape in slide.shapes:
                if shape.has_text_frame and shape != title:
                    shape.text_frame.clear()

            # insert the image to the slide
            left = 0.05
            top = 0.2
            right = left
            
            width_ = 1 - (left + right)
            
            left_ = self.prs.slide_width * left
            top_ = self.prs.slide_height * top
            width_ = self.prs.slide_width * width_
            
            pic = slide.shapes.add_picture(image_path, left_, top_, width=width_)

            if 'note' in toc_item:
                # create a textbox in the slide at the bottom
                left = 0.05
                top = 0.9
                width = 0.9
                height = 0.1

                left_ = self.prs.slide_width * left
                top_ = self.prs.slide_height * top
                width_ = self.prs.slide_width * width
                height_ = self.prs.slide_height * height

                txBox = slide.shapes.add_textbox(left_, top_, width_, height_)
                tf = txBox.text_frame
                tf.text = toc_item['note']
        else:
            slide = self.prs.slides.add_slide(contents_slide_layout)
            title = slide.shapes.title
            title.text = toc_item['title']

    def save(self, filepath):
        """
        Save the presentation

        :param filepath: the file path
        """
        self.filepath = filepath
        self.prs.save(filepath)
    
    def open(self):
        """
        Open the presentation
        """
        if not hasattr(self, "filepath"):
            self.save('/tmp/temp.pptx')
        subprocess.call(['open', self.filepath])